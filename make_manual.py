"""사용자 매뉴얼 PPT 생성 - 모던 디자인"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).parent / "AI세금계산서_사용자매뉴얼.pptx"
APP_NAME = "AI 세금계산서"
APP_TAGLINE = "매입 세금계산서 자동 처리"

# ============================================================
# 디자인 시스템
# ============================================================
# Primary (Indigo)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
PRIMARY_DARK = RGBColor(0x4F, 0x46, 0xE5)
PRIMARY_LIGHT = RGBColor(0xE0, 0xE7, 0xFF)

# Secondary (Emerald)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
SUCCESS_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)

# Accent (Amber)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)

# Danger (Red)
DANGER = RGBColor(0xEF, 0x44, 0x44)

# Neutral grays
TEXT = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
BG_LIGHT = RGBColor(0xF9, 0xFA, 0xFB)
BG_SUBTLE = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "맑은 고딕"


def new_pres():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide, prs, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_rect(slide, x, y, w, h, color, no_line=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if no_line:
        shape.line.fill.background()
    return shape


def add_rounded(slide, x, y, w, h, color, no_line=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if no_line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(0.75)
    return shape


def add_text(slide, x, y, w, h, text, size=14, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_page_header(slide, section_label, page_num, total):
    # 상단 얇은 액센트 스트립
    add_rect(slide, 0, 0, Inches(13.33), Inches(0.08), PRIMARY)

    # 섹션 라벨 (좌측 상단)
    add_text(slide, Inches(0.6), Inches(0.28), Inches(8), Inches(0.3),
             section_label, size=11, color=PRIMARY, bold=True)

    # 페이지 번호 (우측 상단)
    add_text(slide, Inches(11.5), Inches(0.28), Inches(1.5), Inches(0.3),
             f"{page_num:02d} / {total:02d}", size=11, color=MUTED,
             align=PP_ALIGN.RIGHT)


def add_page_footer(slide, prs):
    # 하단 얇은 라인
    add_rect(slide, Inches(0.6), Inches(7.2), Inches(12.13), Emu(9525),
             BORDER)
    add_text(slide, Inches(0.6), Inches(7.25), Inches(8), Inches(0.3),
             f"{APP_NAME}  ·  {APP_TAGLINE}", size=9, color=MUTED)


def add_h1(slide, x, y, w, text):
    # 타이틀 위 액센트 바
    add_rect(slide, x, y, Inches(0.4), Inches(0.05), PRIMARY)
    return add_text(slide, x, y + Inches(0.15), w, Inches(0.7),
                    text, size=34, bold=True, color=TEXT)


def add_subtitle(slide, x, y, w, text):
    return add_text(slide, x, y, w, Inches(0.4),
                    text, size=14, color=MUTED)


# ============================================================
# 슬라이드 생성자들
# ============================================================
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 전체 그라디언트 대신 솔리드 + 액센트
    add_bg(slide, prs, PRIMARY_DARK)
    # 좌측 상단 작은 로고 박스
    add_rounded(slide, Inches(0.8), Inches(0.8), Inches(0.9), Inches(0.9),
                WHITE)
    add_text(slide, Inches(0.8), Inches(0.95), Inches(0.9), Inches(0.6),
             "₩", size=44, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)

    # 앱 이름
    add_text(slide, Inches(0.8), Inches(2.8), Inches(12), Inches(1.5),
             APP_NAME, size=96, bold=True, color=WHITE)

    # 태그라인
    add_text(slide, Inches(0.85), Inches(4.1), Inches(12), Inches(0.8),
             APP_TAGLINE, size=24, color=PRIMARY_LIGHT)

    # 얇은 구분선
    add_rect(slide, Inches(0.8), Inches(5.0), Inches(3), Inches(0.04),
             PRIMARY_LIGHT)

    # 서브
    add_text(slide, Inches(0.8), Inches(5.15), Inches(12), Inches(0.4),
             "사용자 매뉴얼", size=18, color=WHITE)
    add_text(slide, Inches(0.8), Inches(5.6), Inches(12), Inches(0.4),
             "Version 1.0", size=13, color=PRIMARY_LIGHT)

    # 우측 장식 원
    ring = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.8), Inches(4.5), Inches(5), Inches(5)
    )
    ring.fill.background()
    ring.line.color.rgb = PRIMARY_LIGHT
    ring.line.width = Pt(1)
    ring2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.6), Inches(5.3), Inches(3.4), Inches(3.4)
    )
    ring2.fill.background()
    ring2.line.color.rgb = PRIMARY_LIGHT
    ring2.line.width = Pt(0.5)


def slide_toc(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, WHITE)
    add_page_header(slide, "CONTENTS", 2, total)

    add_h1(slide, Inches(0.8), Inches(0.9), Inches(12), "목차")
    add_subtitle(slide, Inches(0.85), Inches(2.0), Inches(12),
                 "이 매뉴얼은 다음과 같이 구성되어 있습니다")

    sections = [
        ("01", "시작하기", "프로그램 소개 · 준비사항 · 첫 실행", PRIMARY),
        ("02", "조회하기", "매입 세금계산서 조회 · 필터링 · PDF 저장", SUCCESS),
        ("03", "매핑 관리", "거래처별 계정/적요 매핑 · 품목 규칙 설정", ACCENT),
        ("04", "GMS 자동화", "전표 입력 · 첨부 · 옵션 · 트러블슈팅", PRIMARY_DARK),
    ]

    base_y = 2.6
    for i, (num, title, desc, color) in enumerate(sections):
        y = Inches(base_y + i * 1.1)
        # 번호 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), y, Inches(0.85), Inches(0.85)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(slide, Inches(0.85), y + Inches(0.12), Inches(0.85),
                 Inches(0.6), num, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

        # 제목
        add_text(slide, Inches(2.0), y + Inches(0.05), Inches(11),
                 Inches(0.5), title, size=22, bold=True, color=TEXT)
        # 설명
        add_text(slide, Inches(2.0), y + Inches(0.5), Inches(11),
                 Inches(0.4), desc, size=13, color=MUTED)

    add_page_footer(slide, prs)


def slide_content(prs, section_label, page_num, total, title,
                  subtitle, bullets, tip=None, two_col=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, WHITE)
    add_page_header(slide, section_label, page_num, total)

    add_h1(slide, Inches(0.8), Inches(0.9), Inches(12), title)
    if subtitle:
        add_subtitle(slide, Inches(0.85), Inches(2.0), Inches(12), subtitle)

    body_top = Inches(2.55)
    if two_col:
        col_w = Inches(5.9)
        col_gap = Inches(0.35)
        left_x = Inches(0.8)
        right_x = left_x + col_w + col_gap
        half = (len(bullets) + 1) // 2
        _draw_bullets(slide, left_x, body_top, col_w, bullets[:half])
        _draw_bullets(slide, right_x, body_top, col_w, bullets[half:])
    else:
        _draw_bullets(slide, Inches(0.8), body_top, Inches(11.7), bullets)

    if tip:
        _draw_tip(slide, Inches(0.8), Inches(6.4), Inches(11.7), tip)

    add_page_footer(slide, prs)


def _draw_bullets(slide, x, y, w, bullets):
    cur_y = y
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if level == 0:
            # 메인 포인트: 작은 점 + 굵은 텍스트
            dot_size = Inches(0.12)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(0.05),
                cur_y + Inches(0.15),
                dot_size, dot_size
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = PRIMARY
            dot.line.fill.background()
            add_text(slide, x + Inches(0.35), cur_y, w - Inches(0.35),
                     Inches(0.4), text, size=14, bold=True, color=TEXT)
            cur_y += Inches(0.42)
        else:
            # 서브 포인트: 들여쓰기 + 중간 텍스트
            add_text(slide, x + Inches(0.7), cur_y, w - Inches(0.7),
                     Inches(0.35), "— " + text, size=12, color=MUTED)
            cur_y += Inches(0.32)


def _draw_tip(slide, x, y, w, text):
    box = add_rounded(slide, x, y, w, Inches(0.6), ACCENT_LIGHT)
    # 좌측 액센트 바
    add_rect(slide, x, y, Inches(0.08), Inches(0.6), ACCENT)
    add_text(slide, x + Inches(0.25), y + Inches(0.15),
             w - Inches(0.3), Inches(0.3),
             "💡  " + text, size=12, color=RGBColor(0x92, 0x40, 0x0E))


def slide_section(prs, section_num, section_label, title, desc,
                  page_num, total, color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, BG_LIGHT)
    add_page_header(slide, section_label, page_num, total)

    # 큰 번호
    add_text(slide, Inches(0.8), Inches(2.2), Inches(12), Inches(1.5),
             section_num, size=140, bold=True, color=PRIMARY_LIGHT)

    # 세로 액센트 바
    add_rect(slide, Inches(0.8), Inches(4.3), Inches(0.08), Inches(1.2), color)

    # 섹션 타이틀
    add_text(slide, Inches(1.0), Inches(4.25), Inches(11), Inches(0.8),
             title, size=40, bold=True, color=TEXT)

    # 설명
    add_text(slide, Inches(1.0), Inches(5.15), Inches(11), Inches(0.6),
             desc, size=16, color=MUTED)

    add_page_footer(slide, prs)


def slide_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs, PRIMARY_DARK)

    add_text(slide, Inches(0.8), Inches(2.8), Inches(12), Inches(1.5),
             "Thank You", size=80, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.8), Inches(4.3), Inches(12), Inches(0.5),
             "문제가 있으시면 개발자에게 문의해주세요",
             size=16, color=PRIMARY_LIGHT, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(5.66), Inches(5.1), Inches(2), Inches(0.04),
             PRIMARY_LIGHT)

    add_text(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(0.5),
             f"{APP_NAME}  ·  v1.0",
             size=13, color=PRIMARY_LIGHT, align=PP_ALIGN.CENTER)


# ============================================================
# 매뉴얼 구성
# ============================================================
def main():
    prs = new_pres()
    total = 22  # 예상 총 페이지 수

    # 1. 표지
    slide_cover(prs)

    # 2. 목차
    slide_toc(prs, total)

    # 3. Section 1 표지
    slide_section(
        prs, "01", "SECTION 01",
        "시작하기",
        "프로그램 설치 없이 exe 하나로 바로 사용할 수 있습니다",
        3, total, PRIMARY,
    )

    # 4. 이 프로그램이 하는 일
    slide_content(
        prs, "Section 01 · 시작하기", 4, total,
        f"{APP_NAME} 이 하는 일",
        "반복되는 매입 세금계산서 업무를 자동화합니다",
        [
            "팝빌 홈택스 매입 세금계산서 조회",
            ("기간 · 거래처 · 과세구분 등 다양한 필터 지원", 1),
            ("한 번 조회로 최대 3개월치 전체 내역", 1),
            "선택 건 PDF 일괄 다운로드",
            ("증빙용 원본 PDF 를 지정한 폴더에 저장", 1),
            "GMS 전표 자동 입력 + PDF 자동 첨부",
            ("거래처 · 계정 · 금액 · 부가세 · 적요 자동 입력", 1),
            ("거래처별 매핑 등록으로 반복 작업 제거", 1),
        ],
        tip="매달 쌓이는 수십 건의 매입 세금계산서 입력을 몇 번의 클릭으로 끝낼 수 있습니다",
    )

    # 5. 준비사항
    slide_content(
        prs, "Section 01 · 시작하기", 5, total,
        "실행 전 준비사항",
        "이 조건만 확인하면 바로 사용할 수 있습니다",
        [
            "Windows 10 또는 11",
            ("Microsoft Edge 기본 설치 (Windows 기본 제공)", 1),
            "인터넷 연결",
            "팝빌 웹사이트에 회사 인증서 등록됨",
            ("이미 회사 계정에 등록되어 있어 별도 작업 불필요", 1),
            "GMS 접속용 O365 계정",
            ("본인 회사 이메일 · 비밀번호 · MFA 준비", 1),
            f"{APP_NAME}.exe 실행 파일",
            ("공유 폴더에서 본인 PC로 복사 후 실행", 1),
        ],
    )

    # 6. 첫 실행
    slide_content(
        prs, "Section 01 · 시작하기", 6, total,
        "프로그램 첫 실행",
        "exe 파일을 더블클릭하면 끝입니다",
        [
            f"{APP_NAME}.exe 더블클릭으로 실행",
            "Windows Defender 경고 시 '추가 정보' → '실행'",
            ("서명이 없는 exe 라 처음에만 경고가 뜰 수 있음", 1),
            "메인 창이 열립니다",
            ("기본 조회 기간이 '이번 달 1일 ~ 오늘' 로 설정됨", 1),
            "처음에는 목록이 비어 있습니다",
            ("오른쪽 상단 '조회' 버튼을 눌러야 목록이 표시됨", 1),
        ],
        tip="최초 실행은 다소 시간이 걸릴 수 있습니다 (내부 리소스 준비)",
    )

    # 7. 메인 화면
    slide_content(
        prs, "Section 01 · 시작하기", 7, total,
        "메인 화면 둘러보기",
        "모든 기능은 이 한 화면에서 이루어집니다",
        [
            "① 조회 조건",
            ("시작일 · 종료일 · 조회 버튼", 1),
            "② 필터 바",
            ("거래처 · 과세 · 유형 · 전표 · 매핑 필터", 1),
            "③ 세금계산서 목록 표",
            ("매핑/전표 표시 · 거래처 · 금액 · 승인번호 등", 1),
            "④ 저장 폴더",
            ("PDF 저장 경로 선택", 1),
            "⑤ GMS 옵션",
            ("백그라운드 실행 · 자동 종료 체크박스", 1),
            "⑥ 액션 버튼",
            ("전체선택 · 거래처 매핑 · 저장(PDF) · GMS 전표 입력", 1),
        ],
    )

    # 8. Section 2 표지
    slide_section(
        prs, "02", "SECTION 02",
        "조회하기",
        "매입 세금계산서 조회 · 필터링 · PDF 저장",
        8, total, SUCCESS,
    )

    # 9. 조회
    slide_content(
        prs, "Section 02 · 조회하기", 9, total,
        "매입 세금계산서 조회",
        "기간을 설정하고 조회 버튼만 누르면 됩니다",
        [
            "시작일 · 종료일 확인 또는 수정",
            ("YYYY-MM-DD 형식 · 최대 3개월 이내", 1),
            "[조회] 버튼 클릭",
            ("내부 처리: 팝빌 → 홈택스 수집 → 목록 반환", 1),
            ("약 10~30초 소요", 1),
            "목록에 전체 매입 세금계산서가 표시됨",
            "동일 기간 재조회는 60분간 캐시 사용",
            ("홈택스 부하 방지를 위함", 1),
        ],
        tip="조회가 끝날 때까지 다른 창에서 자유롭게 작업하실 수 있습니다",
    )

    # 10. 필터
    slide_content(
        prs, "Section 02 · 조회하기", 10, total,
        "필터로 원하는 건만 찾기",
        "5가지 필터를 조합해 정확한 결과를 얻으세요",
        [
            "거래처/사업자번호 검색",
            ("거래처명 또는 사업자번호 일부 입력 → 실시간 필터", 1),
            "과세구분",
            ("전체 · 과세 · 면세 · 영세", 1),
            "문서유형",
            ("전체 · 일반 · 수정 (수정세금계산서 분리)", 1),
            "전표",
            ("전체 · 미입력 · 입력완료 — GMS 입력 여부로 필터", 1),
            "매핑",
            ("전체 · 미매핑 · 매핑완료 — 매핑 등록 여부로 필터", 1),
            "[필터 초기화] 버튼으로 한 번에 리셋",
        ],
        two_col=True,
    )

    # 11. 컬럼 설명
    slide_content(
        prs, "Section 02 · 조회하기", 11, total,
        "목록 표 컬럼",
        "각 컬럼이 의미하는 것",
        [
            "선택 · #",
            ("체크박스 · 순번", 1),
            "매핑 · 전표",
            ("✓ 표시로 매핑/GMS 입력 여부 한눈에 확인", 1),
            "증빙일자 · 발행일자",
            ("세금계산서 작성/발행 날짜", 1),
            "공급자 · 사업자번호",
            ("거래처 정보", 1),
            "유형 · 과세",
            ("일반/수정 · 과세/면세/영세", 1),
            "품목(대표) · 비고",
            ("첫 품목 · 비고 내용", 1),
            "공급가액 · 세액 · 합계",
            ("음수(수정세금계산서)는 빨간색 표시", 1),
            "승인번호",
            ("홈택스 승인번호 · PDF 파일명에 사용됨", 1),
        ],
        two_col=True,
    )

    # 12. PDF 저장
    slide_content(
        prs, "Section 02 · 조회하기", 12, total,
        "PDF 일괄 다운로드",
        "선택한 건을 한 번에 다운로드합니다",
        [
            "저장하려는 건을 체크 (또는 [전체 선택])",
            ("행을 클릭하면 토글 — 선택된 행은 파란색 배경", 1),
            "저장 폴더 확인",
            ("[폴더 선택] 으로 변경 · 설정은 자동 저장", 1),
            ("[열기] 로 탐색기에서 바로 열기 가능", 1),
            "[저장(PDF)] 초록색 버튼 클릭",
            "확인 다이얼로그 → 예",
            "선택 건 수만큼 다운로드 진행",
            ("건당 2~4초 소요 · 진행 상황 하단 상태바 표시", 1),
            "완료 메시지 확인",
        ],
        tip=f"파일명 형식: 승인번호.pdf  —  예: 20260410-42000076-y0100953.pdf",
    )

    # 13. Section 3
    slide_section(
        prs, "03", "SECTION 03",
        "매핑 관리",
        "거래처별 계정 · 적요 · 품목 규칙을 관리합니다",
        13, total, ACCENT,
    )

    # 14. 매핑 개념
    slide_content(
        prs, "Section 03 · 매핑 관리", 14, total,
        "거래처 매핑이란?",
        "GMS 전표 자동 입력을 위한 사전 등록 정보",
        [
            "거래처별로 자동 입력할 값을 미리 등록",
            ("계정 · 코스트센터 · 적요", 1),
            "매핑이 없는 거래처는 GMS 자동 입력 시 자동 스킵",
            ("결과 창에 '매핑 등록 필요 목록' 으로 알려줌", 1),
            "계정/코스트센터는 코드 또는 명 모두 입력 가능",
            ("GMS 에서 Enter 시 자동완성됨", 1),
            "적요는 {월}, {년}, {일} 플레이스홀더 지원",
            ("예: '{월} 유지보수료' → '4월 유지보수료'", 1),
        ],
    )

    # 15. 매핑 등록
    slide_content(
        prs, "Section 03 · 매핑 관리", 15, total,
        "매핑 등록 · 수정",
        "간단한 3단계로 매핑을 관리합니다",
        [
            "대상 선택",
            ("목록에서 거래처 건을 클릭 (자동 프리필)", 1),
            ("또는 매핑 창의 '목록에서 선택' 드롭다운 사용", 1),
            "[거래처 매핑] 버튼 클릭",
            "대화상자에서 값 입력",
            ("계정 · 코스트센터 · 적요", 1),
            ("빈 필드는 GMS 기본값 유지", 1),
            "[거래처+규칙 저장] 클릭",
            "등록된 매핑은 하단 표에 표시",
            ("기존 항목 클릭 → 수정 가능", 1),
        ],
    )

    # 16. 품목 규칙
    slide_content(
        prs, "Section 03 · 매핑 관리", 16, total,
        "품목 규칙 (고급)",
        "같은 거래처라도 품목에 따라 다른 값 적용",
        [
            "같은 거래처의 매입이라도 품목에 따라 구분 필요할 때",
            ("예: '에이시에스' 중 '설치' → 비품, '유지보수' → 지급수수료", 1),
            "대화상자 중간 '품목 규칙' 섹션 사용",
            "품목 포함어 입력 (드롭다운 자동완성 지원)",
            "계정 · 코스트센터 · 적요 입력",
            ("빈 값은 거래처 기본값 상속", 1),
            "[규칙 추가] → 즉시 저장",
            "기존 규칙 클릭 → 폼에 로드 → 수정 후 [규칙 수정]",
            "[규칙 삭제] 로 제거",
        ],
        tip="품목 규칙은 키워드 순서대로 검사됩니다. 먼저 일치한 규칙이 적용됨",
    )

    # 17. Section 4
    slide_section(
        prs, "04", "SECTION 04",
        "GMS 자동화",
        "전표 입력부터 PDF 첨부까지 전 과정 자동",
        17, total, PRIMARY_DARK,
    )

    # 18. 첫 로그인
    slide_content(
        prs, "Section 04 · GMS 자동화", 18, total,
        "GMS 첫 로그인",
        "최초 1회만 수동 로그인하면 됩니다",
        [
            "처음 [GMS 전표 입력] 실행 시 브라우저 창이 열림",
            "'O365 계정으로 로그인' 버튼이 자동 클릭됨",
            "Microsoft 로그인 페이지에서 직접 로그인",
            ("본인 회사 이메일 · 비밀번호", 1),
            ("MFA (폰 승인 등) 완료", 1),
            "로그인 성공 시 세션 자동 저장",
            ("이후부터는 로그인 생략", 1),
            ("세션 파일: %APPDATA%\\PopbillTaxInvoice\\gms_state.json", 1),
            "세션 만료 시 다시 로그인 요청됨",
        ],
        tip="세션은 각자 PC 별로 저장됩니다. 다른 PC 에서 사용 시 재로그인 필요",
    )

    # 19. 전표 입력 실행
    slide_content(
        prs, "Section 04 · GMS 자동화", 19, total,
        "GMS 전표 입력 실행",
        "체크 · 버튼 · 확인 — 3번의 클릭으로 완료",
        [
            "입력할 건들을 체크박스로 선택",
            ("매핑 미등록 건은 자동 스킵됨", 1),
            "GMS 옵션 확인",
            ("백그라운드 실행 · 자동 종료 체크박스", 1),
            "[GMS 전표 입력] 주황 버튼 클릭",
            "확인 창 — 가능 N건 / 스킵 N건 표시",
            ("예 클릭 → 자동 입력 시작", 1),
            "PDF 자동 다운로드 → GMS 자동 입력",
            "완료 후 결과 창",
            ("성공 · 실패 · 스킵 건수 · 매핑 필요 목록", 1),
        ],
    )

    # 20. 옵션 설명
    slide_content(
        prs, "Section 04 · GMS 자동화", 20, total,
        "GMS 옵션",
        "상황에 맞게 옵션을 조합하세요",
        [
            "백그라운드 실행",
            ("체크 시 브라우저 창 없이 처리", 1),
            ("PC 를 자유롭게 사용할 수 있음", 1),
            ("단, 최초 로그인 시에는 자동으로 창이 뜸", 1),
            "완료 후 GMS 창 자동 종료",
            ("체크 시 입력 완료 후 브라우저 자동 닫힘", 1),
            ("체크 해제 시 결과 확인 · 결재 상신 가능", 1),
            "권장 조합",
            ("대량 입력: 백그라운드 + 자동 종료", 1),
            ("검증 필요: 둘 다 해제", 1),
        ],
        two_col=True,
    )

    # 21. FAQ
    slide_content(
        prs, "Section 04 · GMS 자동화", 21, total,
        "자주 묻는 질문",
        "문제가 생기면 여기부터 확인해보세요",
        [
            "Q. 다른 PC 에서도 사용할 수 있나요?",
            ("exe 한 파일만 복사 · 첫 실행 시 MS 로그인 1회", 1),
            "Q. 거래처 매핑은 공유되나요?",
            ("각 PC 별로 저장 · vendor_mapping.json 복사로 공유 가능", 1),
            "Q. 에러가 나면 어디를 봐야 하나요?",
            ("%APPDATA%\\PopbillTaxInvoice\\app.log 확인", 1),
            "Q. 수정세금계산서(음수)도 처리되나요?",
            ("네 · 음수 금액도 자동 입력 · 목록에 빨간색 표시", 1),
            "Q. 이미 입력한 전표를 또 입력하지 않나요?",
            ("매칭 검사로 이미 첨부된 행은 자동 스킵", 1),
        ],
        tip="app.log 는 에러의 상세 정보가 담겨있어 문제 해결에 꼭 필요합니다",
    )

    # 22. End
    slide_end(prs)

    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    main()
